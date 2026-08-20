from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_CONNECTOR
from pathlib import Path

out = Path("/mnt/data/Thermoregulation_Anesthesia_Resident_Seminar.pptx")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Theme
NAVY = RGBColor(17, 38, 63)
TEAL = RGBColor(0, 137, 145)
LIGHT = RGBColor(243, 247, 249)
MID = RGBColor(90, 108, 123)
WHITE = RGBColor(255,255,255)
RED = RGBColor(180, 50, 55)
AMBER = RGBColor(214, 145, 45)
GREEN = RGBColor(55, 125, 88)
BLACK = RGBColor(25, 30, 35)

def add_bg(slide, color=LIGHT):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid(); shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    slide.shapes._spTree.remove(shape._element)
    slide.shapes._spTree.insert(2, shape._element)

def add_header(slide, title, section=None):
    # top bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(.55))
    bar.fill.solid(); bar.fill.fore_color.rgb = NAVY; bar.line.fill.background()
    tx = slide.shapes.add_textbox(Inches(.55), Inches(.08), Inches(10.8), Inches(.35))
    p = tx.text_frame.paragraphs[0]
    p.text = title; p.font.size = Pt(24); p.font.bold = True; p.font.color.rgb = WHITE
    if section:
        st = slide.shapes.add_textbox(Inches(11.2), Inches(.1), Inches(1.55), Inches(.3))
        p = st.text_frame.paragraphs[0]; p.text = section.upper(); p.alignment = PP_ALIGN.RIGHT
        p.font.size = Pt(9); p.font.bold = True; p.font.color.rgb = RGBColor(190,215,220)

def add_footer(slide, num, ref=""):
    # footer line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(.55), Inches(7.14), Inches(12.2), Inches(.02))
    line.fill.solid(); line.fill.fore_color.rgb = TEAL; line.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(.6), Inches(7.17), Inches(11.3), Inches(.22))
    p = tb.text_frame.paragraphs[0]
    p.text = ref[:160]
    p.font.size = Pt(7.5); p.font.color.rgb = MID
    nb = slide.shapes.add_textbox(Inches(12.0), Inches(7.15), Inches(.65), Inches(.25))
    p = nb.text_frame.paragraphs[0]; p.text = str(num); p.alignment = PP_ALIGN.RIGHT
    p.font.size = Pt(8); p.font.bold = True; p.font.color.rgb = MID

def add_bullets(slide, bullets, x=.7, y=1.15, w=7.3, h=5.7, font=18):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.clear(); tf.word_wrap = True
    for i,b in enumerate(bullets):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        if isinstance(b, tuple):
            text, level = b
        else:
            text, level = b, 0
        p.text = text; p.level = level
        p.font.size = Pt(font if level==0 else font-2)
        p.font.color.rgb = BLACK
        p.space_after = Pt(8 if level==0 else 3)
        if level==0:
            p.font.bold = False
    return tb

def add_callout(slide, title, text, x=8.25, y=1.25, w=4.35, h=2.0, fill=WHITE, accent=TEAL):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid(); box.fill.fore_color.rgb = fill
    box.line.color.rgb = accent; box.line.width = Pt(1.5)
    t = slide.shapes.add_textbox(Inches(x+.2), Inches(y+.15), Inches(w-.4), Inches(.35))
    p=t.text_frame.paragraphs[0]; p.text=title; p.font.size=Pt(15); p.font.bold=True; p.font.color.rgb=accent
    b=slide.shapes.add_textbox(Inches(x+.2), Inches(y+.58), Inches(w-.4), Inches(h-.72))
    p=b.text_frame.paragraphs[0]; p.text=text; p.font.size=Pt(12); p.font.color.rgb=BLACK

def add_notes(slide, text):
    tf = slide.notes_slide.notes_text_frame
    tf.text = text

def add_diagram_box(slide, text, x, y, w, h, fill=WHITE, accent=TEAL, font=13):
    sh=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb=fill
    sh.line.color.rgb=accent; sh.line.width=Pt(1.3)
    tf=sh.text_frame; tf.clear(); tf.vertical_anchor=MSO_ANCHOR.MIDDLE
    p=tf.paragraphs[0]; p.text=text; p.alignment=PP_ALIGN.CENTER
    p.font.size=Pt(font); p.font.bold=True; p.font.color.rgb=NAVY
    return sh

def arrow(slide, x1,y1,x2,y2, color=TEAL):
    ln=slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    ln.line.color.rgb=color; ln.line.width=Pt(2)
    ln.line.end_arrowhead = True
    return ln

def add_title_slide():
    s=prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s,NAVY)
    # accent
    a=s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(6.75), prs.slide_width, Inches(.75))
    a.fill.solid(); a.fill.fore_color.rgb=TEAL; a.line.fill.background()
    t=s.shapes.add_textbox(Inches(.75), Inches(1.35), Inches(11.8), Inches(1.4))
    p=t.text_frame.paragraphs[0]; p.text="THERMOREGULATION\nAND ITS APPLICATION IN ANESTHESIA"
    p.font.size=Pt(34); p.font.bold=True; p.font.color.rgb=WHITE
    st=s.shapes.add_textbox(Inches(.78), Inches(3.05), Inches(10.5), Inches(.55))
    p=st.text_frame.paragraphs[0]; p.text="From Physiology to Perioperative Temperature Management"
    p.font.size=Pt(20); p.font.color.rgb=RGBColor(185,220,224)
    info=s.shapes.add_textbox(Inches(.8), Inches(4.35), Inches(7.0), Inches(1.35))
    tf=info.text_frame; tf.clear()
    for i,txt in enumerate(["Resident: __________________________",
                            "Department of Anaesthesiology: __________________________",
                            "Institution: __________________________",
                            "Date: __________________________"]):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.text=txt; p.font.size=Pt(14); p.font.color.rgb=WHITE; p.space_after=Pt(7)
    add_notes(s, "Opening: Introduce thermoregulation as an active physiological process rather than simply a temperature reading. Use the opening case: a 55-year-old patient undergoing a 4-hour abdominal procedure develops a core temperature of 34.2°C despite a seemingly normal OR environment. Ask: what happened physiologically, why did anesthesia cause it, what complications can follow, and how should the anesthesiologist respond?")
    return s

add_title_slide()

slides = []

def make_slide(title, bullets, notes, ref="", section="", callout=None, diagram=None):
    s=prs.slides.add_slide(prs.slide_layouts[6]); add_bg(s); add_header(s,title,section)
    if bullets: add_bullets(s,bullets)
    if callout: add_callout(s,*callout)
    if diagram: diagram(s)
    add_footer(s,len(prs.slides),ref)
    add_notes(s,notes)
    return s

# 2
make_slide("Opening Case: The 34.2°C Patient",
["55-year-old; 4-hour abdominal surgery under general anesthesia",
 "Core temperature falls to 34.2°C despite apparently normal OR conditions",
 "Ask: redistribution, impaired defenses, environmental heat loss — or all three?",
 "Potential consequences: bleeding, arrhythmia, drug prolongation, delayed recovery, shivering",
 "The anesthesiologist must prevent, detect, explain and treat temperature disturbances"],
"Use the case to frame the seminar. The key concept is that anesthesia widens the range over which core temperature can drift before thermoregulatory responses are triggered. Induction commonly causes vasodilation and redistribution of heat from the core to peripheral tissues. Ongoing heat loss then continues through radiation, convection, conduction and evaporation. The management question is not merely 'how to warm' but 'why did the patient cool, how severe is it, and what complications are emerging?'",
"Source brief supplied by user; perioperative framework aligned with NICE CG65.","CLINICAL HOOK",
("Resident question","Do not wait for severe hypothermia before acting: temperature trends matter clinically."))

make_slide("Learning Objectives",
["Explain normal human thermoregulation and the inter-threshold range",
 "Describe how general, neuraxial and intravenous anesthetics alter thermal control",
 "Recognize redistribution hypothermia and its three intraoperative phases",
 "Choose appropriate temperature-monitoring sites and technologies",
 "Prevent and treat perioperative hypothermia using active warming and fluid warming",
 "Differentiate perioperative hyperthermia and manage malignant hyperthermia promptly",
 "Apply principles to trauma, pediatrics, obstetrics, ICU and postoperative care"],
"By the end, the resident should be able to connect physiology to bedside decisions. Emphasize that temperature management is part of routine anesthetic safety, not an optional add-on.",
"User-provided presentation brief, lines 1030–1050.","FOUNDATION")

def thermo_diagram(s):
    add_diagram_box(s,"Thermal sensors",.7,2.1,2.1,1.0)
    add_diagram_box(s,"Spinal cord\n+ brainstem",3.3,2.1,2.1,1.0)
    add_diagram_box(s,"Hypothalamus\nIntegration",5.9,2.1,2.1,1.0)
    add_diagram_box(s,"Effectors",8.5,2.1,1.8,1.0)
    add_diagram_box(s,"Vasoconstrict\nShiver • Sweat\nBehavior",10.8,1.65,1.7,1.9,fill=RGBColor(238,248,249))
    arrow(s,2.8,2.6,3.3,2.6); arrow(s,5.4,2.6,5.9,2.6); arrow(s,8.0,2.6,8.5,2.6); arrow(s,10.3,2.6,10.8,2.6)
    add_callout(s,"Core principle","Thermoregulation is a feedback system: detect → integrate → compare with regulated range → recruit effectors.",.9,4.4,11.3,1.4,fill=WHITE,accent=TEAL)

make_slide("Normal Thermoregulation: The Control System",
["Core temperature is regulated within a relatively narrow physiological range",
 "Thermal information arises from peripheral and central thermoreceptors",
 "The hypothalamus integrates temperature signals against defended thresholds",
 "Responses include vasoconstriction, shivering, sweating and non-shivering thermogenesis",
 "Behavioral responses are powerful but are largely unavailable during anesthesia"],
"Explain the control loop. The hypothalamus is not simply a thermometer; it coordinates autonomic, somatic and behavioral responses. The distinction between core and peripheral temperature is crucial because induction changes blood flow between compartments.",
"Guyton & Hall; Sessler DI. Temperature monitoring and perioperative thermoregulation. Anesthesiology. 2008;109:318–338.","PHYSIOLOGY",diagram=thermo_diagram)

make_slide("Core, Peripheral and Mean Body Temperature",
["Core compartment: thoracic/abdominal/cranial organs; most clinically relevant for anesthetic management",
 "Peripheral tissues act as a thermal reservoir and can be substantially cooler than the core",
 "Core-to-peripheral gradient is dynamic and changes with vasoconstriction or vasodilation",
 "Mean body temperature reflects the combined thermal state of core and peripheral compartments",
 "A peripheral fall can precede or accompany core changes; after induction, redistribution narrows the gradient rapidly"],
"Stress that 'body temperature' is not a single homogeneous value. Peripheral temperature can be low while core temperature is initially preserved. After induction, vasodilation mixes cooler peripheral blood with the core, producing the classic early temperature drop.",
"Sessler 2008; user-provided brief lines 55–81.","PHYSIOLOGY")

def threshold_diagram(s):
    # axes
    ln=s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1.2), Inches(5.8), Inches(11.8), Inches(5.8)); ln.line.color.rgb=MID
    ln=s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(2.0), Inches(6.1), Inches(2.0), Inches(1.4)); ln.line.color.rgb=MID
    # bands
    add_diagram_box(s,"Cold-defense zone\nVasoconstriction → shivering",1.8,1.8,3.3,1.1,fill=RGBColor(238,246,252),accent=TEAL)
    add_diagram_box(s,"INTER-THRESHOLD RANGE\nLittle autonomic response",5.0,3.1,3.3,1.1,fill=WHITE,accent=GREEN)
    add_diagram_box(s,"Heat-defense zone\nVasodilation → sweating",8.2,1.8,3.3,1.1,fill=RGBColor(255,246,238),accent=AMBER)
    arrow(s,5.1,2.35,5.0,3.1); arrow(s,8.2,3.65,8.2,2.9)
    add_callout(s,"Under anesthesia","Volatile and IV agents lower the thresholds for vasoconstriction and shivering, widening the range in which temperature can drift.",2.0,5.0,9.3,1.0,fill=WHITE,accent=RED)

make_slide("Thermoregulatory Thresholds & Inter-threshold Range",
["Cold and warm defenses are triggered at thresholds rather than at one fixed 'set point'",
 "The inter-threshold range is the interval between vasoconstriction and sweating thresholds",
 "Cold-defense hierarchy: vasoconstriction → shivering → increased metabolic thermogenesis",
 "Heat-defense hierarchy: vasodilation → sweating; behavioral cooling normally adds major protection",
 "Anesthetic drugs shift thresholds downward, particularly for vasoconstriction and shivering"],
"Sessler's threshold model is high-yield for viva questions. The inter-threshold range is narrow in awake humans but becomes wider under anesthesia because drug effects suppress autonomic responses. Exact threshold values vary with age, sex, circadian rhythm, drugs and clinical state; avoid presenting one number as universal.",
"Sessler 2008; user-provided brief lines 232–249.","PHYSIOLOGY",diagram=threshold_diagram)

make_slide("Thermal Inputs: Peripheral & Central Thermoreceptors",
["Skin contains cold- and warm-sensitive receptors that provide anticipatory information",
 "Central thermosensitive regions include the preoptic/anterior hypothalamus and spinal cord",
 "Blood temperature provides information from deep thoracic and abdominal structures",
 "Cold signals are especially important for initiating vasoconstriction and shivering",
 "Central and peripheral inputs are integrated before an effector response is recruited"],
"Explain that peripheral thermal input is not merely redundant. Skin temperature can influence thresholds and helps the CNS anticipate changes in core temperature. Deep-body temperature signals provide information about the thermal state of critical organs.",
"Guyton & Hall; Sessler 2008.","PHYSIOLOGY")

def heatloss(s):
    labels=[("RADIATION",1.0,2.0,RED),("CONVECTION",4.1,2.0,TEAL),("CONDUCTION",7.2,2.0,AMBER),("EVAPORATION",10.3,2.0,GREEN)]
    desc=["Infrared exchange with surroundings","Airflow removes heat","Direct contact with cold surfaces/fluids","Liquid-to-vapor heat loss"]
    for (lab,x,y,c),d in zip(labels,desc):
        add_diagram_box(s,lab,x,y,2.1,1.0,accent=c)
        tx=s.shapes.add_textbox(Inches(x),Inches(3.2),Inches(2.1),Inches(1.0))
        p=tx.text_frame.paragraphs[0]; p.text=d; p.alignment=PP_ALIGN.CENTER; p.font.size=Pt(11); p.font.color.rgb=MID
    add_callout(s,"OR relevance","Large exposed surface + low ambient temperature + cold fluids + open cavities can make heat loss clinically important.",1.1,4.8,11.1,1.15,fill=WHITE,accent=TEAL)

make_slide("Heat Exchange: The Four Mechanisms",
["Radiation: net infrared transfer between the patient and surrounding surfaces",
 "Convection: heat carried away by moving air; increased by airflow",
 "Conduction: direct transfer to tables, mattresses, irrigation fluid or other surfaces",
 "Evaporation: latent heat loss from skin, open cavities, preparation fluids and the respiratory tract",
 "In the OR, mechanisms occur simultaneously; prevention therefore requires layered interventions"],
"Do not overemphasize a single percentage contribution because it varies with patient, operation and environment. Forced-air warming works primarily by transferring heat through warmed moving air across a large surface area.",
"User-provided brief lines 137–173; NICE CG65.","PHYSIOLOGY",diagram=heatloss)

make_slide("Cold Responses: Vasoconstriction, Shivering & Thermogenesis",
["Vasoconstriction is an early sympathetic response that reduces cutaneous blood flow and conserves core heat",
 "Shivering is involuntary skeletal-muscle activity that can markedly increase metabolic heat production",
 "Shivering raises oxygen consumption, carbon dioxide production and cardiovascular workload",
 "Non-shivering thermogenesis relies on brown adipose tissue and mitochondrial uncoupling via UCP-1",
 "Neonates and infants depend more heavily on non-shivering thermogenesis because shivering capacity is immature"],
"Shivering is not a benign symptom. It can increase oxygen demand and interfere with monitoring, wound comfort and recovery. In neonates, brown adipose tissue is a major heat-producing organ. UCP-1 uncouples oxidative phosphorylation so energy is released as heat.",
"Guyton & Hall; user-provided brief lines 177–216.","PHYSIOLOGY")

make_slide("Heat Responses",
["Peripheral vasodilation increases cutaneous blood flow and heat transfer to the environment",
 "Sweating increases evaporative heat loss when ambient conditions permit evaporation",
 "Behavioral cooling is normally a major component of human heat defense",
 "Anesthesia blunts behavioral responses and alters autonomic thresholds",
 "Hyperthermia management depends on identifying the cause rather than treating the temperature number alone"],
"Contrast physiologic heat defense with perioperative hyperthermia. Sweating is often ineffective in a humid environment and can be absent in anticholinergic toxicity. In malignant hyperthermia, excessive heat is generated by skeletal muscle hypermetabolism; the key treatment is dantrolene and removal of triggers.",
"Sessler 2008; user-provided brief lines 218–229.","PHYSIOLOGY")

make_slide("How General Anesthesia Disrupts Thermoregulation",
["Reduces the thresholds for vasoconstriction and shivering",
 "Produces peripheral vasodilation and increases redistribution of core heat",
 "Suppresses behavioral responses to cold and heat",
 "Neuromuscular blockade prevents effective shivering despite metabolic demand",
 "Depth of anesthesia, agent selection, regional techniques and environmental exposure all influence risk",
 "The result is a larger temperature drift with less physiological compensation"],
"Emphasize the interaction of mechanisms. Anesthetic-induced vasodilation causes redistribution, while suppressed thresholds reduce the ability to defend temperature. Neuromuscular blockade removes shivering as an effective heat-producing response. This explains why a patient can become hypothermic even in an OR that does not feel cold.",
"NICE CG65; Sessler 2008; user-provided brief lines 252–274.","ANESTHESIA")

make_slide("Anesthetic Agents: Practical Thermal Effects",
["Volatile anesthetics: dose-dependent reduction in vasoconstriction and shivering thresholds",
 "Propofol and other IV anesthetics also impair autonomic thermoregulation; induction commonly promotes redistribution",
 "Opioids reduce shivering thresholds and can contribute to postoperative shivering control at selected doses",
 "Alpha-2 agonists can reduce shivering but may cause bradycardia, hypotension and sedation",
 "Neuromuscular blockers abolish effective shivering but do not restore thermoregulatory thresholds",
 "Clinical effect is usually drug + dose + technique + environment, not one agent in isolation"],
"Use this as a comparative rather than a memorization slide. Avoid implying that every listed drug has an identical thermal effect. The clinically important point is that anesthetic depth and combinations change thresholds and heat redistribution. Drug effects are particularly relevant during long cases, neuraxial anesthesia and emergence.",
"Sessler 2008; user-provided brief lines 276–316.","ANESTHESIA")

make_slide("Neuraxial Anesthesia: Why an Awake Patient Can Become Hypothermic",
["Spinal and epidural anesthesia block sympathetic vasoconstriction below the level of block",
 "Regional vasodilation redistributes heat toward blocked peripheral tissues",
 "Thermal perception is altered and the patient may not recognize the degree of cooling",
 "Shivering can be suppressed or displaced by the block",
 "Patients can therefore become hypothermic despite being awake and conversational",
 "High-risk settings include cesarean delivery, prolonged orthopedic and lower abdominal surgery"],
"Explain the paradox: consciousness does not guarantee intact thermoregulation. Neuraxial blockade interrupts autonomic and afferent pathways and can impair both vasoconstriction and shivering. In obstetric anesthesia, maternal hypothermia may be overlooked because attention is focused on blood pressure, nausea, neuraxial level and neonatal considerations.",
"User-provided brief lines 318–342.","ANESTHESIA")

def phases(s):
    # line chart
    x0,y0=1.0,5.9; x1,y1=11.8,1.4
    ln=s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x0), Inches(y0), Inches(11.8), Inches(y0)); ln.line.color.rgb=MID
    ln=s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x0), Inches(y0), Inches(x0), Inches(y1)); ln.line.color.rgb=MID
    pts=[(1.2,2.0),(3.5,3.1),(6.8,4.7),(10.0,4.95),(11.5,4.95)]
    for a,b in zip(pts,pts[1:]):
        ln=s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(a[0]), Inches(a[1]), Inches(b[0]), Inches(b[1])); ln.line.color.rgb=TEAL; ln.line.width=Pt(3)
    for x,label in [(2.2,"Redistribution"),(5.2,"Linear heat loss"),(9.2,"Plateau")]:
        tb=s.shapes.add_textbox(Inches(x),Inches(6.05),Inches(2.0),Inches(.35)); p=tb.text_frame.paragraphs[0]; p.text=label; p.font.size=Pt(11); p.font.bold=True; p.font.color.rgb=NAVY
    add_callout(s,"Temperature trend","The initial drop is usually fastest after induction; later, heat loss slows as vasoconstriction/plateau develops.",2.2,1.1,8.6,1.0,fill=WHITE,accent=TEAL)

make_slide("Perioperative Hypothermia: The Three Classic Phases",
["Phase 1 — Redistribution: rapid core temperature fall after induction",
 "Phase 2 — Linear heat loss: continued loss through radiation, convection, conduction and evaporation",
 "Phase 3 — Plateau: vasoconstriction reduces net heat loss; temperature stabilizes if heat production approximates loss",
 "The plateau may be absent when vasoconstriction is strongly suppressed or heat loss is extreme",
 "Recognizing the phase helps the anesthesiologist choose prevention versus rescue strategies"],
"Use the curve to explain why simply turning on a warmer late may not immediately reverse the initial fall. Prewarming expands peripheral heat content before induction and reduces the magnitude of redistribution. During phase 2, active warming and exposure control are central.",
"NICE CG65; Sessler 2008; user-provided brief lines 344–397.","ANESTHESIA",diagram=phases)

make_slide("Risk Factors for Perioperative Hypothermia",
["Patient: older age, neonates, low body mass, frailty, pre-existing hypothermia, trauma and sepsis",
 "Surgery: long duration, large exposed surface, open cavities, major blood loss",
 "Anesthesia: general or neuraxial anesthesia, deep anesthesia, vasodilation and neuromuscular blockade",
 "Environment: cool OR, cold irrigation, unwarmed IV fluids/blood and excessive skin exposure",
 "High-risk patients benefit from a proactive warming plan before induction rather than rescue alone"],
"Turn risk assessment into a pre-induction checklist. The combination of patient vulnerability and procedure duration is more important than any one factor. Trauma and sepsis add impaired perfusion and altered metabolic states.",
"User-provided brief lines 399–440.","PREVENTION")

make_slide("Clinical Consequences of Perioperative Hypothermia",
["Cardiovascular: sympathetic activation, hypertension/vasoconstriction, arrhythmias and myocardial stress",
 "Coagulation: platelet dysfunction and impaired coagulation enzyme activity → increased bleeding risk",
 "Drug effects: slower metabolism/clearance and prolonged neuromuscular blockade",
 "Respiratory/metabolic: increased oxygen demand and carbon dioxide production during shivering",
 "Immune and wound effects: impaired tissue oxygenation and potential increase in surgical-site infection risk",
 "Recovery: shivering, discomfort, delayed emergence and prolonged PACU/hospital care"],
"Keep causal language appropriately cautious: perioperative hypothermia is associated with several adverse outcomes, and mechanisms are biologically plausible. Avoid claiming every outcome is caused solely by temperature. The key clinical message is that even mild hypothermia can be consequential.",
"NICE CG65; user-provided brief lines 442–508.","CONSEQUENCES")

make_slide("Perioperative Hyperthermia: A Diagnostic Problem",
["Infection/sepsis and transfusion reactions can raise temperature",
 "Drug-related syndromes include anticholinergic toxicity, serotonin syndrome and neuroleptic malignant syndrome",
 "Endocrine causes include thyroid storm",
 "Malignant hyperthermia is a hypermetabolic skeletal-muscle crisis triggered by susceptible patients exposed to triggering anesthetics",
 "Environmental heat and iatrogenic warming can contribute",
 "Always interpret temperature with ETCO₂, heart rate, rigidity, acid-base status, potassium and clinical context"],
"Hyperthermia is a sign, not a diagnosis. In an anesthetized patient, rapidly rising ETCO2 with tachycardia and rigidity should trigger immediate consideration of malignant hyperthermia even before marked hyperthermia develops.",
"User-provided brief lines 512–529.","HYPERTHERMIA")

def mh(s):
    add_diagram_box(s,"Trigger\nVolatile anesthetic / succinylcholine",.6,1.55,2.5,1.0,accent=RED)
    add_diagram_box(s,"RYR1 / Ca²⁺\nrelease",3.65,1.55,2.1,1.0,accent=RED)
    add_diagram_box(s,"Muscle\nhypermetabolism",6.25,1.55,2.1,1.0,accent=RED)
    add_diagram_box(s,"↑ ETCO₂ • acidosis\nrigidity • K⁺ • heat",8.85,1.55,3.3,1.0,accent=RED)
    arrow(s,3.1,2.05,3.65,2.05,RED); arrow(s,5.75,2.05,6.25,2.05,RED); arrow(s,8.35,2.05,8.85,2.05,RED)
    add_callout(s,"Key point","Hyperthermia may be a late sign. Early clues are unexplained rising ETCO₂, tachycardia, rigidity and metabolic acidosis.",1.2,3.2,10.9,1.35,fill=WHITE,accent=RED)

make_slide("Malignant Hyperthermia: Pathophysiology",
["MH is an inherited susceptibility of skeletal-muscle calcium handling; RYR1 is a major implicated gene",
 "Triggering agents include volatile anesthetics and succinylcholine",
 "Abnormal calcium release causes sustained muscle contraction and hypermetabolism",
 "Early manifestations may include rapidly rising ETCO₂, tachycardia, rigidity and acidosis",
 "Hyperkalemia, rhabdomyolysis, hyperthermia and myoglobinuria can follow",
 "Hyperthermia is important — but often later than the metabolic warning signs"],
"Stress the diagnostic pattern. Do not wait for temperature to become extreme. The combination of rapidly increasing ETCO2 despite adequate ventilation, tachycardia, rigidity and metabolic/respiratory acidosis is highly concerning. RYR1 is the major gene, but MH susceptibility is genetically heterogeneous.",
"MHAUS; user-provided brief lines 531–568.","EMERGENCY",diagram=mh)

make_slide("Malignant Hyperthermia: Immediate Management Algorithm",
["1. Recognize MH and call for help / activate the MH protocol",
 "2. Stop all triggering agents; use non-triggering anesthesia",
 "3. Hyperventilate with 100% oxygen; use high fresh-gas flow and charcoal filters if available",
 "4. Give IV dantrolene 2.5 mg/kg rapidly; repeat until the hypermetabolic response improves",
 "5. Treat hyperkalemia, severe acidosis, arrhythmias and rhabdomyolysis; monitor urine output",
 "6. Begin active cooling when indicated; stop cooling around 38°C to avoid overshoot",
 "7. Transfer to ICU for continued monitoring and recurrence surveillance"],
"Current MHAUS guidance supports an initial dantrolene dose of 2.5 mg/kg based on true body weight, repeated as needed. MHAUS notes that doses >10 mg/kg may sometimes be required in persistent crises. Prioritize dantrolene, trigger removal, hyperventilation and metabolic management; cooling is an adjunct, not a substitute for dantrolene.",
"MHAUS crisis guidance and recommendations.","EMERGENCY")

make_slide("Temperature Monitoring: What Counts as 'Core'?",
["Pulmonary artery: reference standard in invasive monitoring, but rarely justified solely for temperature",
 "Distal esophagus: useful during general anesthesia with an airway device; responds relatively rapidly to core changes",
 "Nasopharynx: useful when probe position is appropriate and brain/core temperature is the target",
 "Tympanic/aural methods: convenient but technique-dependent and can be inaccurate for core temperature",
 "Bladder and rectal sites: slower response; affected by urine flow and local conditions",
 "Site selection must match the clinical question, procedure and expected temperature dynamics"],
"Explain that no site is universally 'best'. Distal esophageal and nasopharyngeal sites are commonly used during general anesthesia, while invasive pulmonary artery temperature is highly accurate when a pulmonary artery catheter is already indicated. Aural infrared measurements can be convenient but are not interchangeable with continuous core monitoring.",
"Association of Anaesthetists monitoring guidance; NICE CG65.","MONITORING")

make_slide("Temperature Measurement Technologies",
["Thermistors: resistance changes with temperature; common in disposable probes",
 "Thermocouples: voltage generated across dissimilar metals; robust and widely used",
 "Infrared thermometry: detects emitted infrared radiation; convenient but sensitive to site and technique",
 "Zero-heat-flux systems: create an insulated skin surface and estimate core temperature non-invasively",
 "Disposable probes are practical, but placement and calibration still determine accuracy",
 "Trend consistency and clinically appropriate site often matter as much as nominal device accuracy"],
"Describe the technology at a level useful for a resident: what the sensor measures, how quickly it responds, and where errors occur. Zero-heat-flux systems are attractive when invasive or mucosal probes are impractical, but local validation and device-specific limitations matter.",
"User-provided brief lines 603–615; current device literature should be checked locally before procurement.","MONITORING")

make_slide("Evidence-Based Monitoring: Practical Standards",
["NICE CG65: measure and document temperature before induction and every 30 min intraoperatively; postoperative monitoring on recovery admission and at intervals",
 "NICE defines inadvertent perioperative hypothermia as core temperature <36.0°C",
 "NICE recommends warming patients undergoing anesthesia >30 min from induction; higher-risk short cases also warrant active warming",
 "Association of Anaesthetists standards include temperature in minimum monitoring for procedures >30 min and throughout recovery",
 "For long/complex cases, continuous core-temperature monitoring is preferable where feasible",
 "Local institutional policy may be more specific than these baseline standards"],
"Important nuance: NICE CG65 is an older guideline and its exact implementation may be updated locally. Present it as a foundational guideline, not as the only current authority. The Association of Anaesthetists monitoring standard states temperature for procedures over 30 minutes and includes temperature during recovery.",
"NICE CG65; Association of Anaesthetists 2016 monitoring guidance.","MONITORING")

make_slide("Prevention: The Perioperative Temperature Bundle",
["Preoperative: measure temperature, identify risk and avoid unnecessary exposure",
 "Prewarm selected patients before induction to increase peripheral heat content",
 "Prepare active warming equipment before induction rather than after temperature falls",
 "Minimize exposed surface area while maintaining surgical access",
 "Warm IV fluids, blood products and irrigation fluids when clinically indicated",
 "Continue warming and temperature monitoring into PACU until temperature is acceptable"],
"Prevention works best as a bundle. The objective is to reduce redistribution, reduce ongoing environmental loss and replace heat continuously. Do not rely on blankets alone for patients at meaningful risk.",
"NICE CG65; user-provided brief lines 617–649.","PREVENTION")

make_slide("Active vs Passive Warming",
["Active: forced-air, resistive/conductive systems and radiant warming deliver heat to the patient",
 "Passive: blankets, drapes, insulation and head covering reduce heat loss",
 "Passive measures are useful adjuncts but generally cannot replace active warming during significant heat loss",
 "Forced-air warming has a strong evidence base for prevention and treatment of perioperative hypothermia",
 "Choice depends on surgical access, equipment, patient factors and local infection/fire-safety protocols"],
"Forced-air warming remains a widely used standard because it can deliver substantial heat over a large surface. Resistive systems are alternatives, particularly when forced-air access is limited. Device-specific safety instructions and local infection-control policy must be followed.",
"NICE CG65; user-provided brief lines 671–686.","PREVENTION")

make_slide("Fluid, Blood & Irrigation Warming",
["Cold IV fluids can directly reduce body temperature, especially at high infusion rates",
 "Warm larger-volume IV fluids and blood products with validated fluid/blood warming systems",
 "Warming is particularly important during massive transfusion, trauma, prolonged procedures and pediatric surgery",
 "NICE CG65 recommends warming IV fluids ≥500 mL and blood products to about 37°C",
 "Avoid improvised warming methods that can cause overheating, hemolysis or inaccurate temperature delivery"],
"Do not place blood products in unvalidated warming devices. Use equipment intended for blood/fluid warming. The NICE recommendation is specifically about 500 mL or more of IV fluid and blood products, but clinical judgment is needed for lower volumes administered rapidly.",
"NICE CG65; user-provided brief lines 651–669.","PREVENTION")

make_slide("Prewarming: Prevent Redistribution Before It Starts",
["Prewarming transfers heat into peripheral tissues before induction",
 "A warmer peripheral compartment reduces the core-to-peripheral gradient after anesthetic vasodilation",
 "The practical goal is to reduce the magnitude of redistribution hypothermia",
 "Use active warming while maintaining patient comfort and monitoring",
 "For a 3-hour laparoscopic case: assess risk → prewarm → prepare active intraoperative warming → warm fluids → monitor core temperature"],
"Explain the mechanism rather than presenting a single universal prewarming duration. Evidence supports preoperative active warming, but protocols vary by device and patient. In a viva, say that prewarming should begin sufficiently before induction to meaningfully increase peripheral heat content, commonly around 10–30 minutes depending on the device and workflow.",
"NICE CG65 update evidence; user-provided brief lines 688–702.","PREVENTION")

make_slide("Intraoperative Temperature Management Algorithm",
["BEFORE INDUCTION → measure temperature → risk-stratify → prewarm if indicated → prepare warmer + warmed fluids",
 "AFTER INDUCTION → establish appropriate core monitoring → minimize exposure → start active warming",
 "DURING SURGERY → trend temperature → warm fluids/blood → adjust warming intensity → reassess causes of unexpected change",
 "IF <36°C → continue/strengthen active warming, investigate contributing factors and monitor for complications",
 "POSTOPERATIVE → continue warming, treat shivering, document temperature and ensure acceptable recovery temperature"],
"Use this as the resident's bedside algorithm. A low temperature should trigger both treatment and a search for cause: phase of anesthesia, warming failure, excessive exposure, cold fluids, hemorrhage, sepsis or endocrine issues.",
"NICE CG65; user-provided brief lines 704–737.","PREVENTION")

make_slide("Postoperative Shivering",
["Causes: hypothermia, redistribution, neuraxial anesthesia, pain and altered thermoregulatory thresholds",
 "Consequences: discomfort, increased oxygen consumption and carbon dioxide production, cardiovascular stress",
 "First steps: confirm temperature, provide active warming, reduce exposure and address pain",
 "Pharmacologic options include pethidine/meperidine, clonidine, dexmedetomidine and tramadol",
 "Drug selection depends on hemodynamics, respiratory risk, sedation and local practice",
 "Differentiate shivering from seizure, rigors, anxiety, pain and drug-related movement"],
"Drug evidence varies and local formularies differ. Pethidine/meperidine has historically been used for post-anesthetic shivering but can cause sedation, respiratory depression and nausea. Alpha-2 agonists can reduce shivering but may cause hypotension/bradycardia. Always treat the underlying thermal problem.",
"User-provided brief lines 739–765; pharmacology evidence should be interpreted in local context.","RECOVERY")

make_slide("Special Populations: Pediatrics, Geriatrics & Obstetrics",
["Neonates/infants: high surface-area-to-volume ratio, immature thermoregulation and reliance on brown fat",
 "Children: rapid heat loss during exposure, fluid administration and prolonged surgery",
 "Older adults: reduced vasoconstrictor and shivering responses and less thermoregulatory reserve",
 "Obstetric anesthesia: neuraxial blockade can produce maternal hypothermia and shivering",
 "Cesarean delivery: consider maternal warming, fluid warming and neonatal temperature implications",
 "Risk assessment should be population-specific rather than simply age-based"],
"Highlight that pediatric warming must avoid both hypothermia and thermal injury. In older adults, modest temperature changes may produce disproportionate consequences. During cesarean section, neuraxial anesthesia does not eliminate hypothermia risk.",
"User-provided brief lines 769–799.","SPECIAL POPULATIONS")

make_slide("Trauma, Massive Transfusion & the Lethal Cycle",
["Trauma patients lose heat through exposure, hemorrhage, cold environments and administration of unwarmed fluids/blood",
 "Hypothermia impairs coagulation and worsens bleeding",
 "Acidosis further reduces coagulation efficiency and myocardial performance",
 "Transfusion of large volumes of cold blood can deepen hypothermia",
 "Core cycle: Hypothermia → coagulopathy → bleeding → transfusion → further hypothermia",
 "Management: rapid hemorrhage control + warmed blood/fluids + active external warming + temperature monitoring"],
"Frame hypothermia as a component of trauma physiology rather than a secondary nuisance. Damage-control resuscitation should be paired with aggressive temperature management. Avoid stating a single universal trauma temperature target unless tied to a specific protocol.",
"User-provided brief lines 847–866.","SPECIAL POPULATIONS")

make_slide("Neurosurgery, Cardiac Surgery & Intentional Hypothermia",
["Cardiac surgery may deliberately alter temperature during cardiopulmonary bypass and requires controlled rewarming",
 "Rewarming must avoid excessive gradients and rapid temperature overshoot",
 "Neurosurgery may use temperature manipulation for selected neuroprotection strategies",
 "Intentional hypothermia is different from inadvertent perioperative hypothermia",
 "Targeted temperature management after cardiac arrest is now framed around controlled temperature strategies rather than reflexive deep hypothermia",
 "Evidence and targets are evolving; follow contemporary ICU/cardiac-arrest protocols"],
"Distinguish three concepts: accidental hypothermia, perioperative inadvertent hypothermia, and deliberate targeted temperature management. The latter is disease-specific and protocol-driven. Avoid transplanting a cardiac-arrest target into routine anesthesia.",
"User-provided brief lines 828–845.","SPECIAL POPULATIONS")

make_slide("ICU Temperature Management",
["Fever and hypothermia are physiologic signals requiring contextual interpretation",
 "Continuous temperature monitoring is appropriate in unstable or neurologically vulnerable patients",
 "Sepsis: temperature is one component of a broader perfusion and infection assessment",
 "Post-cardiac-arrest patients may require protocolized targeted temperature management",
 "Neurological patients may be particularly sensitive to temperature excursions",
 "Treat the underlying cause while using controlled temperature management when indicated"],
"Keep ICU management separate from routine perioperative warming. Fever can be adaptive or harmful depending on context, and aggressive temperature control has trade-offs. Use disease-specific protocols.",
"User-provided brief lines 868–880.","ICU")

make_slide("Differential Diagnosis of Intraoperative Hyperthermia",
["Malignant hyperthermia: rapid ETCO₂ rise, rigidity, tachycardia, acidosis, hyperkalemia",
 "Serotonin syndrome: serotonergic exposure + clonus/hyperreflexia, GI/autonomic features",
 "Neuroleptic malignant syndrome: dopamine antagonist exposure + rigidity, altered mental state, slower evolution",
 "Anticholinergic toxicity: dry skin, mydriasis, tachycardia, reduced sweating",
 "Thyroid storm: hypermetabolism, tachyarrhythmia, CNS/GI features and known thyroid disease",
 "Sepsis/transfusion/environmental causes: assess timing, exposures, cultures, hemolysis and hemodynamics"],
"Use the pattern of onset and associated findings. In MH, ETCO2 rise and hypermetabolism are often early. Serotonin syndrome is characterized by neuromuscular hyperactivity such as clonus. NMS typically evolves more slowly. Anticholinergic toxicity produces impaired heat dissipation through anhidrosis.",
"User-provided brief lines 512–529; MHAUS for MH.","HYPERTHERMIA")

make_slide("Clinical Case 1: Elective Surgery with Hypothermia",
["3-hour laparoscopic abdominal case; induction temperature 36.5°C",
 "After 40 min: 35.1°C; forced-air warmer was not started",
 "Question: phase of hypothermia?",
 "Interpretation: redistribution followed by ongoing heat loss",
 "Management: active warming, minimize exposure, warm fluids, monitor trend and complications",
 "Prevention: prewarming + active warming from induction + core monitoring"],
"Expected answer: predominantly phase 1 transitioning into phase 2. The case illustrates why waiting for a temperature below 36°C is inferior to proactive prevention. Ask the resident what caused the first fall and what could have been done before induction.",
"NICE CG65; user-provided clinical case framework.","CASES")

make_slide("Clinical Case 2: Shivering after Spinal Anesthesia",
["Cesarean section under spinal anesthesia; patient is awake but shivering",
 "Temperature is 35.4°C; patient reports feeling 'cold' but is hemodynamically stable",
 "Consider neuraxial vasodilation + impaired thermoregulation + redistribution",
 "Immediate: warming, blankets/forced-air where feasible, warm IV fluids",
 "If persistent and clinically significant: consider pharmacologic treatment after assessing contraindications",
 "Continue maternal and neonatal temperature surveillance"],
"Discuss that shivering after spinal anesthesia is common and multifactorial. Avoid reflexively treating the movement without treating hypothermia. Drug choice should be individualized and consistent with obstetric practice.",
"User-provided brief lines 739–765 and 789–799.","CASES")

make_slide("Clinical Case 3: Trauma + Hypothermia + Coagulopathy",
["Polytrauma patient arrives cold after prolonged exposure; temperature 34.0°C",
 "Ongoing hemorrhage; fibrinogen/platelets falling; acidosis present",
 "Recognize the hypothermia–coagulopathy–bleeding cycle",
 "Immediate priorities: hemorrhage control, damage-control resuscitation and aggressive warming",
 "Warm blood/fluids, use active external warming and minimize exposure",
 "Monitor temperature continuously and reassess coagulation/acid-base status"],
"Expected answer: treat temperature as part of the resuscitation problem. The patient's coagulopathy is not solved by warming alone; hemorrhage control and balanced resuscitation are essential. The cycle can accelerate quickly.",
"User-provided brief lines 801–810 and 847–866.","CASES")

make_slide("Clinical Case 4: Suspected Malignant Hyperthermia",
["Sevoflurane anesthesia; sudden ETCO₂ rise despite increased minute ventilation",
 "Tachycardia + generalized rigidity + mixed respiratory/metabolic acidosis",
 "Temperature is initially only 37.8°C",
 "Diagnosis: treat as suspected MH — do not wait for severe hyperthermia",
 "Stop triggers, 100% oxygen/hyperventilation, dantrolene 2.5 mg/kg IV and repeat as needed",
 "Treat hyperkalemia/acidosis, monitor CK/urine output and transfer to ICU"],
"This is a classic viva trap: temperature is not required to make the diagnosis early. The early metabolic signal is often unexplained ETCO2 elevation. MHAUS recommends prompt dantrolene and trigger discontinuation.",
"MHAUS crisis management.","CASES")

make_slide("Clinical Case 5: Pediatric Prolonged Surgery",
["Infant undergoing prolonged surgery with large exposed surface area",
 "Temperature drops quickly despite normal adult-sized OR practices",
 "Risk: high surface-area-to-volume ratio + immature thermoregulation + limited heat reserve",
 "Use appropriately sized active warming, warmed fluids and minimized exposure",
 "Avoid thermal injury: monitor temperature continuously and follow device instructions",
 "Postoperative warming and temperature surveillance should continue until stable"],
"Discuss why adult protocols cannot simply be scaled down. Pediatric patients can change temperature rapidly, and both hypothermia and overheating are dangerous. Equipment must be pediatric-appropriate.",
"User-provided brief lines 773–781.","CASES")

make_slide("High-Yield Summary Table: Prevention & Rescue",
["PREVENT → measure temp → risk assess → prewarm → minimize exposure → active warming",
 "WARM → forced-air/resistive/radiant + warmed IV fluids/blood/irrigation when indicated",
 "MONITOR → choose appropriate core site; trend rather than rely on isolated readings",
 "RESCUE → if hypothermic, intensify warming and identify ongoing causes",
 "RECOVER → continue warming and treat shivering; document temperature",
 "EMERGENCY → hyperthermia + hypermetabolism = consider MH and act immediately"],
"Use this as the final clinical framework before the viva section. It is intentionally algorithmic: the resident should be able to recite it during ward/theatre rounds.",
"NICE CG65; Association of Anaesthetists; MHAUS.","SUMMARY")

make_slide("30 High-Yield Viva Questions — Part 1",
["1. What is thermoregulation?  2. Define core temperature.  3. What is the inter-threshold range?",
 "4. Name four mechanisms of heat loss.  5. Why does temperature fall rapidly after induction?",
 "6. What is redistribution hypothermia?  7. Why can spinal anesthesia cause hypothermia?",
 "8. What are the major cold-defense mechanisms?  9. Why is shivering clinically important?",
 "10. What is non-shivering thermogenesis?  11. Why are neonates at high risk?",
 "12. What temperature defines perioperative hypothermia in NICE CG65?  13. Name three consequences of hypothermia.",
 "14. Which temperature sites are commonly used during GA?  15. Why can tympanic infrared readings be unreliable?"],
"Model answers: 1) Coordinated neural/autonomic/behavioral control of body temperature. 2) Temperature of well-perfused deep tissues/organs. 3) Interval between vasoconstriction and sweating thresholds. 4) Radiation, convection, conduction, evaporation. 5) Anesthetic vasodilation mixes cooler peripheral blood into the core and suppresses defenses. 6) Core heat redistribution to peripheral tissues after induction. 7) Sympathetic blockade causes vasodilation and impairs vasoconstriction/shivering. 8) Vasoconstriction, shivering, non-shivering thermogenesis. 9) It increases metabolic and cardiopulmonary demand. 10) Brown-fat thermogenesis via UCP-1. 11) High surface-area-to-volume ratio and immature responses. 12) Core <36°C. 13) Bleeding, arrhythmias, delayed drug clearance, shivering, delayed recovery. 14) Distal esophagus, nasopharynx, rectal, bladder, tympanic and invasive sites depending on context. 15) Technique/site limitations and poor agreement with core temperature.",
"User-provided viva requirements; NICE CG65; Association of Anaesthetists.","VIVA")

make_slide("30 High-Yield Viva Questions — Part 2",
["16. How does hypothermia affect coagulation?  17. How does it affect drug metabolism?",
 "18. What are the phases of intraoperative hypothermia?",
 "19. What is prewarming and why does it work?",
 "20. Compare active and passive warming.  21. Why warm blood and IV fluids?",
 "22. What are causes of postoperative shivering?  23. Name pharmacologic options.",
 "24. What is the earliest important clue to MH?  25. Why can hyperthermia be late in MH?",
 "26. What triggers MH?  27. What is the initial dantrolene dose?",
 "28. How do you differentiate MH from serotonin syndrome/NMS?  29. What is the trauma lethal cycle?  30. What is targeted temperature management?"],
"Model answers: 16) Hypothermia impairs platelet function and coagulation enzyme kinetics. 17) Clearance/metabolism may slow, prolonging drug effects. 18) Redistribution, linear heat loss, plateau. 19) Active warming before induction increases peripheral heat content and reduces redistribution. 20) Active adds heat; passive reduces heat loss. 21) Large volumes of cold fluid/blood can cause substantial heat loss. 22) Hypothermia, neuraxial anesthesia, pain and altered thresholds. 23) Pethidine/meperidine, clonidine, dexmedetomidine, tramadol among studied agents. 24) Unexplained rapid ETCO2 rise with tachycardia/rigidity. 25) It is a consequence of sustained hypermetabolism and may follow earlier metabolic signs. 26) Volatile anesthetics and succinylcholine. 27) 2.5 mg/kg IV based on true weight, repeated as needed. 28) Use exposure history + neuromuscular pattern + tempo + ETCO2/metabolic features. 29) Hypothermia→coagulopathy→bleeding→transfusion→further hypothermia. 30) Protocolized temperature control for selected critical illnesses, distinct from accidental hypothermia.",
"MHAUS; NICE CG65; user-provided viva requirements.","VIVA")

make_slide("20 Difficult Faculty Questions — Part 1",
["1. Why can a patient be cold peripherally while core temperature is initially normal?",
 "2. Why is prewarming physiologically different from simply warming after induction?",
 "3. Why is '36°C' useful clinically but not a universal biological set point?",
 "4. Why might neuraxial patients become hypothermic despite being awake?",
 "5. Why can neuromuscular blockade worsen the clinical consequences of cooling?",
 "6. Why is a single temperature reading less informative than a trend?",
 "7. What determines whether a site is a good proxy for core temperature?",
 "8. Why is forced-air warming effective, and when is it difficult to use?"],
"Model answers: 1) Peripheral tissues can be a cooler reservoir while vasoconstriction limits mixing. 2) Prewarming increases peripheral heat content before vasodilation; post-induction warming alone may lag behind redistribution. 3) 36°C is a pragmatic perioperative threshold associated with increased risk below it, not a fixed hypothalamic set point. 4) Sympathetic blockade and altered thermal perception impair defenses. 5) It prevents effective shivering. 6) Trends reveal phase, rate of change and response to interventions. 7) Anatomic proximity to central blood flow, perfusion, site-specific artifact and response time. 8) It transfers heat over a broad surface but may conflict with surgical access, device positioning or local infection/fire protocols.",
"Sessler 2008; NICE CG65; Association of Anaesthetists.","FACULTY")

make_slide("20 Difficult Faculty Questions — Part 2",
["9. Why is ETCO₂ more useful than temperature as an early MH clue?",
 "10. When should active cooling be emphasized in MH?",
 "11. Why is dantrolene the specific therapy for MH?",
 "12. Why can very large dantrolene doses be required?",
 "13. What is the difference between MH and sepsis-related hyperthermia?",
 "14. How does hypothermia worsen trauma coagulopathy?",
 "15. What are limitations of NICE CG65 when applied to modern practice?",
 "16. What is the role of zero-heat-flux monitoring?",
 "17. Why should cooling in MH not distract from dantrolene and metabolic treatment?",
 "18. What are the major pediatric warming hazards?",
 "19. Why is intentional hypothermia different from inadvertent hypothermia?",
 "20. What is an important evidence gap in perioperative temperature management?"],
"Model answers: 9) ETCO2 reflects rapidly increasing CO2 production from skeletal-muscle hypermetabolism and can rise before temperature. 10) When core temperature is significantly elevated or rising, while simultaneously treating the hypermetabolic crisis. 11) It reduces skeletal-muscle calcium release and hypermetabolism. 12) Severe/persistent crises may require repeated dosing. 13) MH is exposure-triggered hypermetabolism with characteristic ETCO2/rigidity pattern; sepsis usually has infectious/inflammatory context and different tempo. 14) Platelet and enzymatic coagulation function are temperature dependent. 15) It is an older guideline; local updated evidence/device standards may supersede individual recommendations. 16) It offers non-invasive estimation of core temperature but requires device-specific validation. 17) Because ongoing hypermetabolism is the cause of heat generation. 18) Burns, overheating, rapid temperature shifts and inappropriate adult equipment. 19) Intentional strategies are controlled, indication-specific and protocolized. 20) Optimal prewarming duration, best non-invasive monitoring technology and comparative effectiveness of warming systems remain areas where practice/device evidence continues to evolve.",
"MHAUS; NICE CG65; Association of Anaesthetists; user-provided faculty section.","FACULTY")

make_slide("Key Clinical Algorithms — One-Slide Takeaway",
["HYPOTHERMIA: Measure → prevent redistribution → active warming → warm fluids → monitor trend → treat shivering",
 "UNEXPECTED TEMPERATURE FALL: verify measurement → check probe/site → assess exposure/fluids/warming device → assess hemorrhage/sepsis",
 "HYPERTHERMIA: check ETCO₂ + rigidity + drugs + neuromuscular signs → consider MH early",
 "SUSPECTED MH: stop triggers → 100% O₂/hyperventilate → dantrolene → metabolic/arrhythmia treatment → cooling when indicated → ICU",
 "TRAUMA: hemorrhage control + warmed blood/fluids + active warming + serial temperature/coagulation/acid-base assessment"],
"Use this slide to rehearse the seminar closing. The resident should be able to move from temperature measurement to a cause-directed algorithm within seconds.",
"NICE CG65; MHAUS.","ALGORITHMS")

make_slide("Take-Home Messages",
["Temperature is a physiological variable, not merely another monitor number",
 "Anesthesia lowers thermoregulatory thresholds and promotes redistribution of core heat",
 "Prewarming + active intraoperative warming + warmed fluids is more effective than rescue alone",
 "Monitor temperature with a site/device appropriate to the procedure and expected dynamics",
 "Even mild perioperative hypothermia can have clinically relevant consequences",
 "Hyperthermia requires diagnosis; suspected MH demands immediate action before severe hyperthermia develops",
 "The best temperature strategy is proactive, continuous, and integrated with the entire perioperative pathway"],
"Close by returning to the opening case. A 34.2°C patient is not simply 'cold': the temperature reflects disrupted physiology with downstream consequences. The anesthesiologist can intervene at every step — before induction, during surgery, during emergence and in recovery.",
"NICE CG65; MHAUS; Sessler 2008.","CLOSING")

make_slide("Selected Bibliography & Verified Guidance",
["Sessler DI. Temperature monitoring and perioperative thermoregulation. Anesthesiology. 2008;109:318–338.",
 "NICE. Inadvertent perioperative hypothermia: The management of inadvertent perioperative hypothermia in adults undergoing surgery (CG65).",
 "Checketts MR, et al. Recommendations for standards of monitoring during anaesthesia and recovery 2015. Anaesthesia. 2016;71:85–93.",
 "MHAUS. Managing a Crisis / recommendations on dantrolene and MH treatment.",
 "Miller’s Anesthesia; Barash: Clinical Anesthesia; Morgan & Mikhail’s Clinical Anesthesiology.",
 "Guyton & Hall Textbook of Medical Physiology; Ganong’s Review of Medical Physiology."],
"Use the bibliography as a starting reference list. For a departmental submission, the resident should verify the latest institutional policy and any newer society guidance before presenting. The source brief explicitly requested that references be real and verifiable.",
"Selected sources verified during preparation; additional textbook editions should be checked against the resident's institution.","REFERENCES")

# add a final "reference notes" to last slide with web verification details
add_notes(prs.slides[-1], 
"""Verified web sources used to update the deck:
• NICE CG65: temperature before induction, every 30 min intraoperatively, postoperative monitoring; hypothermia <36°C; warming recommendations.
• Association of Anaesthetists monitoring guidance: temperature as minimum monitoring for procedures >30 min and recovery; continuous core monitoring recommended for longer/complex cases in commentary.
• MHAUS: suspected MH management, stop triggers, 100% oxygen/hyperventilation, dantrolene 2.5 mg/kg IV, repeat as needed, active cooling when indicated.
The deck intentionally labels NICE CG65 as foundational/older guidance rather than claiming it is the newest universal standard."""
)

# Add a small visual reference slide for "suggested diagrams"
make_slide("Suggested Diagrams for Faculty Discussion",
["1. Thermoregulatory control loop: sensors → spinal/brainstem → hypothalamus → effectors",
 "2. Core/peripheral compartments and redistribution after induction",
 "3. Heat loss mechanisms: radiation, convection, conduction, evaporation",
 "4. Threshold graph: vasoconstriction, shivering, sweating and inter-threshold range",
 "5. Three-phase intraoperative temperature curve",
 "6. MH calcium-release pathway and emergency management algorithm",
 "7. Trauma hypothermia–coagulopathy vicious cycle"],
"These are already represented as editable vector diagrams in the deck where practical. They can be replaced with institutional figures if desired.",
"User-provided brief lines 938–950.","DIAGRAMS")

# Put bibliography slide before diagrams? Keep current order okay.
# Save
prs.save(out)
print(f"Created {out} with {len(prs.slides)} slides.")
